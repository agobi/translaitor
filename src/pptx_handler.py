"""PowerPoint (PPTX) document handler for extraction and reintegration."""

from typing import Any

from pptx import Presentation

from .base_handler import BaseDocumentHandler


def extract_text_from_shape(shape):
    """Extract text from a shape, extracting each RUN separately to preserve formatting.

    Args:
        shape: A shape object from python-pptx

    Returns:
        list: List of text strings (one per run with formatting)
    """
    texts = []

    # Handle shapes with text frames - extract each RUN separately
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:  # Include runs with any text
                    texts.append(run.text)

    # Handle tables - extract each cell's runs
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            texts.append(run.text)

    # Handle grouped shapes (recursively)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))

    return texts


class TextIterator:
    """Iterator for translated texts to match extraction order."""

    def __init__(self, translated_data):
        self.slides = translated_data["slides"]
        self.current_slide = 0
        self.current_text = 0

    def get_next(self):
        """Get next translated text."""
        if self.current_slide >= len(self.slides):
            return None

        slide_texts = self.slides[self.current_slide]["texts"]

        if self.current_text >= len(slide_texts):
            return None

        text = slide_texts[self.current_text]
        self.current_text += 1
        return text

    def next_slide(self):
        """Move to next slide."""
        self.current_slide += 1
        self.current_text = 0


def reintegrate_text_into_shape(shape, text_iterator):
    """Reintegrate translated text into a shape.

    Args:
        shape: A shape object from python-pptx
        text_iterator: Iterator for translated texts

    Returns:
        int: Number of text elements (runs) replaced
    """
    count = 0

    # Handle shapes with text frames - replace each RUN
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    translated_text = text_iterator.get_next()
                    if translated_text is not None:
                        run.text = translated_text
                        count += 1

    # Handle tables - replace each cell's runs
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            translated_text = text_iterator.get_next()
                            if translated_text is not None:
                                run.text = translated_text
                                count += 1

    # Handle grouped shapes (recursively)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            count += reintegrate_text_into_shape(sub_shape, text_iterator)

    return count


class PPTXHandler(BaseDocumentHandler):
    """Handler for PowerPoint (PPTX) files."""

    def extract_text(self, pptx_path):
        """Extract all text from a PPTX file in deterministic order.

        Args:
            pptx_path: Path to the PPTX file

        Returns:
            dict: JSON structure with slides and texts
        """
        presentation = Presentation(pptx_path)

        result: dict[str, list[dict[str, Any]]] = {"slides": []}

        for slide in presentation.slides:
            slide_texts = []

            # Iterate through shapes in order
            for shape in slide.shapes:
                texts = extract_text_from_shape(shape)
                slide_texts.extend(texts)

            # Always add slide entry, even if empty
            result["slides"].append({"texts": slide_texts})

        return result

    def print_extraction_summary(self, data):
        """Print extraction summary.

        Args:
            data: Extracted data dictionary
        """
        total_texts = sum(len(slide["texts"]) for slide in data["slides"])
        print(f"✓ Extracted {total_texts} text elements from {len(data['slides'])} slides")

    def reintegrate_text(self, pptx_path, translated_data, output_path):
        """Reintegrate translated text back into PPTX.

        Args:
            pptx_path: Path to original PPTX file
            translated_data: Dictionary with translated texts
            output_path: Path to save translated PPTX

        Returns:
            int: Total number of text elements replaced
        """
        presentation = Presentation(pptx_path)
        text_iterator = TextIterator(translated_data)

        total_replaced = 0

        for slide in presentation.slides:
            # Process shapes in same order as extraction
            for shape in slide.shapes:
                count = reintegrate_text_into_shape(shape, text_iterator)
                total_replaced += count

            # Move to next slide in iterator
            text_iterator.next_slide()

        # Save the modified presentation
        presentation.save(output_path)

        return total_replaced
