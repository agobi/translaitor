"""Extract text from PPTX files to JSON format."""

import json
from pptx import Presentation


def extract_text_from_shape(shape):
    """Extract text from a shape, handling various shape types.
    
    Args:
        shape: A shape object from python-pptx
        
    Returns:
        list: List of text strings found in the shape
    """
    texts = []
    
    # Handle shapes with text frames
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            texts.append(text)
    
    # Handle tables
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text_frame.text.strip()
                if text:
                    texts.append(text)
    
    # Handle grouped shapes (recursively)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))
    
    return texts


def extract_text_from_pptx(pptx_path):
    """Extract all text from a PPTX file in deterministic order.
    
    Args:
        pptx_path: Path to the PPTX file
        
    Returns:
        dict: JSON structure with slides and texts
    """
    presentation = Presentation(pptx_path)
    
    result = {"slides": []}
    
    for slide in presentation.slides:
        slide_texts = []
        
        # Iterate through shapes in order
        for shape in slide.shapes:
            texts = extract_text_from_shape(shape)
            slide_texts.extend(texts)
        
        # Always add slide entry, even if empty
        result["slides"].append({"texts": slide_texts})
    
    return result


def save_extracted_json(data, output_path):
    """Save extracted data to JSON file.
    
    Args:
        data: Dictionary with extracted data
        output_path: Path to save JSON file
    """
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def extract(pptx_path, json_path):
    """Main extraction function.
    
    Args:
        pptx_path: Path to input PPTX file
        json_path: Path to output JSON file
    """
    data = extract_text_from_pptx(pptx_path)
    save_extracted_json(data, json_path)
    
    # Count total texts
    total_texts = sum(len(slide["texts"]) for slide in data["slides"])
    print(f"✓ Extracted {total_texts} text elements from {len(data['slides'])} slides")
    print(f"✓ Saved to {json_path}")
