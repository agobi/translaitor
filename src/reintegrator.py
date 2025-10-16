"""Reintegrate translated text back into PPTX files."""

import json
from pptx import Presentation


class TextIterator:
    """Iterator for translated texts to match extraction order."""
    
    def __init__(self, translated_data):
        self.slides = translated_data["slides"]
        self.current_slide = 0
        self.current_text = 0
    
    def get_next(self):
        """Get next translated text."""
        if self.current_slide >= len(self.slides):
            return None
        
        slide_texts = self.slides[self.current_slide]["texts"]
        
        if self.current_text >= len(slide_texts):
            return None
        
        text = slide_texts[self.current_text]
        self.current_text += 1
        return text
    
    def next_slide(self):
        """Move to next slide."""
        self.current_slide += 1
        self.current_text = 0


def set_shape_text(shape, text):
    """Set text in a shape, handling various shape types.
    
    Args:
        shape: A shape object from python-pptx
        text: Text to set
        
    Returns:
        int: Number of text elements replaced
    """
    count = 0
    
    # Handle shapes with text frames
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        if shape.text_frame.text.strip():
            shape.text_frame.text = text
            count = 1
    
    return count


def replace_text_preserve_formatting(text_frame, new_text):
    """Replace text while preserving formatting from the first run.
    
    Strategy: Save formatting, clear all text, put new text in first run with saved formatting.
    
    Args:
        text_frame: Text frame to modify
        new_text: New text to set
    """
    if not text_frame.paragraphs:
        text_frame.text = new_text
        return
    
    # Get the first paragraph
    first_paragraph = text_frame.paragraphs[0]
    
    if not first_paragraph.runs:
        text_frame.text = new_text
        return
    
    # Get the first run and save its font properties
    first_run = first_paragraph.runs[0]
    
    # Save font properties
    saved_font = {
        'name': None,
        'size': None,
        'bold': None,
        'italic': None,
        'underline': None,
        'color_rgb': None
    }
    
    try:
        saved_font['name'] = first_run.font.name
    except (AttributeError, KeyError):
        pass
    
    try:
        saved_font['size'] = first_run.font.size
    except (AttributeError, KeyError):
        pass
    
    try:
        saved_font['bold'] = first_run.font.bold
    except (AttributeError, KeyError):
        pass
    
    try:
        saved_font['italic'] = first_run.font.italic
    except (AttributeError, KeyError):
        pass
    
    try:
        saved_font['underline'] = first_run.font.underline
    except (AttributeError, KeyError):
        pass
    
    try:
        if first_run.font.color and first_run.font.color.type is not None:
            saved_font['color_rgb'] = first_run.font.color.rgb
    except (AttributeError, KeyError):
        pass
    
    # Clear all text using the simple method
    text_frame.text = new_text
    
    # Now apply the saved formatting to the first run
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        new_first_run = text_frame.paragraphs[0].runs[0]
        
        try:
            if saved_font['name']:
                new_first_run.font.name = saved_font['name']
        except:
            pass
        
        try:
            if saved_font['size']:
                new_first_run.font.size = saved_font['size']
        except:
            pass
        
        try:
            if saved_font['bold'] is not None:
                new_first_run.font.bold = saved_font['bold']
        except:
            pass
        
        try:
            if saved_font['italic'] is not None:
                new_first_run.font.italic = saved_font['italic']
        except:
            pass
        
        try:
            if saved_font['underline'] is not None:
                new_first_run.font.underline = saved_font['underline']
        except:
            pass
        
        try:
            if saved_font['color_rgb']:
                new_first_run.font.color.rgb = saved_font['color_rgb']
        except:
            pass


def reintegrate_text_into_shape(shape, text_iterator):
    """Reintegrate translated text into a shape.
    
    Args:
        shape: A shape object from python-pptx
        text_iterator: Iterator for translated texts
        
    Returns:
        int: Number of text elements replaced
    """
    count = 0
    
    # Handle shapes with text frames
    if hasattr(shape, "text_frame") and shape.has_text_frame:
        if shape.text_frame.text.strip():
            translated_text = text_iterator.get_next()
            if translated_text is not None:
                replace_text_preserve_formatting(shape.text_frame, translated_text)
                count += 1
    
    # Handle tables
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame.text.strip():
                    translated_text = text_iterator.get_next()
                    if translated_text is not None:
                        replace_text_preserve_formatting(cell.text_frame, translated_text)
                        count += 1
    
    # Handle grouped shapes (recursively)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            count += reintegrate_text_into_shape(sub_shape, text_iterator)
    
    return count


def reintegrate_text_into_pptx(pptx_path, translated_data, output_path):
    """Reintegrate translated text back into PPTX.
    
    Args:
        pptx_path: Path to original PPTX file
        translated_data: Dictionary with translated texts
        output_path: Path to save translated PPTX
    """
    presentation = Presentation(pptx_path)
    text_iterator = TextIterator(translated_data)
    
    total_replaced = 0
    
    for slide_idx, slide in enumerate(presentation.slides):
        # Process shapes in same order as extraction
        for shape in slide.shapes:
            count = reintegrate_text_into_shape(shape, text_iterator)
            total_replaced += count
        
        # Move to next slide in iterator
        text_iterator.next_slide()
    
    # Save the modified presentation
    presentation.save(output_path)
    
    return total_replaced


def load_translated_json(json_path):
    """Load translated JSON file.
    
    Args:
        json_path: Path to JSON file
        
    Returns:
        dict: Translated data
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def reintegrate(original_pptx_path, translated_json_path, output_pptx_path):
    """Main reintegration function.
    
    Args:
        original_pptx_path: Path to original PPTX file
        translated_json_path: Path to translated JSON file
        output_pptx_path: Path to save translated PPTX
    """
    translated_data = load_translated_json(translated_json_path)
    
    total_replaced = reintegrate_text_into_pptx(
        original_pptx_path,
        translated_data,
        output_pptx_path
    )
    
    total_expected = sum(len(slide["texts"]) for slide in translated_data["slides"])
    
    print(f"✓ Replaced {total_replaced} text elements")
    
    if total_replaced != total_expected:
        print(f"⚠ Warning: Expected {total_expected} texts but replaced {total_replaced}")
    
    print(f"✓ Saved to {output_pptx_path}")
